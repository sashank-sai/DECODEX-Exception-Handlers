"""
DECODE X 2026 — Stage 2 Interim Submission: 5-Slide PPTX Builder
==================================================================
Generates the mandatory 5-slide structural recalibration brief.
Convert to PDF before submission.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

DATA_DIR = Path(r'c:\Users\asus\Desktop\decodex')
CHART_DIR = DATA_DIR / 'charts' / 'stage2'
OUT_FILE = DATA_DIR / 'submission' / 'Stage2_Recalibration_Brief.pptx'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x0D, 0x47, 0xA1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x26, 0x32, 0x38)
GRAY = RGBColor(0x75, 0x75, 0x75)
RED = RGBColor(0xD3, 0x2F, 0x2F)
GREEN = RGBColor(0x38, 0x8E, 0x3C)
ORANGE = RGBColor(0xF5, 0x7C, 0x00)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)

def add_bg(slide, color=LIGHT_BG):
    """Set slide background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title_text, subtitle_text=''):
    """Add a dark blue title bar at the top."""
    from pptx.util import Inches, Pt
    from pptx.oxml.ns import qn
    # Title bar shape
    left, top, width, height = Inches(0), Inches(0), Inches(13.333), Inches(1.1)
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.15)
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
        p2.font.italic = True

def add_textbox(slide, left, top, width, height, text, font_size=11, bold=False, color=DARK, alignment=PP_ALIGN.LEFT):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf

def add_rich_textbox(slide, left, top, width, height, lines):
    """Add textbox with multiple styled lines. lines = [(text, size, bold, color), ...]"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return tf

def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    """Add a formatted table."""
    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    table = table_shape.table
    
    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.alignment = PP_ALIGN.CENTER
                # Color negative values red, positive green
                if isinstance(val, str) and val.startswith('-'):
                    p.font.color.rgb = RED
                elif isinstance(val, str) and val.startswith('+'):
                    p.font.color.rgb = GREEN
    
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)
    
    return table

def add_chart_image(slide, filename, left, top, width):
    """Add a chart image."""
    path = CHART_DIR / filename
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width))
    else:
        # Try stage1 charts folder
        alt_path = DATA_DIR / 'charts' / 'stage1' / filename
        if alt_path.exists():
            slide.shapes.add_picture(str(alt_path), Inches(left), Inches(top), Inches(width))

def add_callout(slide, left, top, width, height, text, bg_color=RGBColor(0xFF, 0xF9, 0xC4), border_color=ORANGE):
    """Add a highlighted callout box."""
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = DARK

# ============================================================
# SLIDE 1: SHOCK DIAGNOSIS
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
add_bg(slide1)
add_title_bar(slide1, 'SLIDE 1: Shock Diagnosis', 'Metro Phase 2 — Structural Break Detection')

# Chart: Forecast deviation waterfall
add_chart_image(slide1, 's2_01_forecast_deviation.png', 0.3, 1.3, 7)

# Right panel: Key findings
add_rich_textbox(slide1, 7.5, 1.3, 5.5, 1.5, [
    ('STRUCTURAL BREAK: +16.1%', 16, True, RED),
    ('Q3 actual exceeded Stage 1 forecast by 402,832 passengers', 11, False, DARK),
    ('Classification: REGIME CHANGE (not temporary shock)', 11, True, DARK),
])

# Shift classification table
add_table(slide1, 7.5, 3.0, 5.5, 1.8,
    ['Shift Type', 'Finding', 'Magnitude'],
    [
        ['Level', 'Express ↑, Feeder ↓', '+29-32% / -18-31%'],
        ['Volatility', 'All routes LOWER CV', '-2 to -18%'],
        ['Elasticity', 'Congestion sensitivity 2×', '+101%'],
        ['Congestion', 'Higher despite metro', '+10.9%'],
    ])

# Route deviation table
add_table(slide1, 7.5, 5.0, 5.5, 2.0,
    ['Route', 'Type', 'Deviation'],
    [
        ['X66', 'Express', '+51.3%'],
        ['X11', 'Express', '+49.3%'],
        ['X28', 'Express', '+47.1%'],
        ['C04', 'City', '+20.5%'],
        ['F18', 'Feeder', '-22.2%'],
        ['F12', 'Feeder', '-8.9%'],
    ])

# ============================================================
# SLIDE 2: MODEL RECALIBRATION SUMMARY
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2)
add_title_bar(slide2, 'SLIDE 2: Model Recalibration Summary', 'Stage 1 → Stage 2 Forecast Revision')

# Chart: Timeline with break
add_chart_image(slide2, 's2_03_timeline_break.png', 0.3, 1.3, 7)

# Method explanation
add_rich_textbox(slide2, 7.5, 1.3, 5.5, 1.2, [
    ('Recalibration Methodology', 14, True, DARK_BLUE),
    ('1. Compute Q3 forecast-vs-actual deviation per route', 10, False, DARK),
    ('2. Classify: Stable (<5%) | Moderate (5-15%) | Regime (>15%)', 10, False, DARK),
    ('3. Apply adjustment: 75% permanent for regime, 50% for moderate', 10, False, DARK),
])

# Revised forecast table
add_table(slide2, 7.5, 2.8, 5.5, 1.4,
    ['Metric', 'Stage 1', 'Stage 2', 'Change'],
    [
        ['Q4 Total', '2,999,357', '3,366,963', '+12.3%'],
        ['Uncertainty', '±5%', '±23.6%', '4.7× wider'],
        ['Full-Year 2025', '11,142,733', '11,910,399', '+6.9%'],
    ])

# Route-level recalibration
add_table(slide2, 7.5, 4.5, 5.5, 2.3,
    ['Route', 'Type', 'Classification', 'Factor', 'Q4 Change'],
    [
        ['X66', 'Express', 'Regime Change', '1.385', '+38.5%'],
        ['X11', 'Express', 'Regime Change', '1.370', '+37.0%'],
        ['X28', 'Express', 'Regime Change', '1.353', '+35.3%'],
        ['C04', 'City', 'Regime Change', '1.154', '+15.4%'],
        ['F18', 'Feeder', 'Regime Change', '0.834', '-16.6%'],
        ['F12', 'Feeder', 'Moderate Shift', '0.956', '-4.4%'],
    ])

add_callout(slide2, 0.3, 6.5, 7, 0.7,
    '⚡ Key: We separate temporary shock (25%) from permanent regime change (75%). '
    'The wider ±23.6% uncertainty reflects limited post-shock data (1 quarter only).')

# ============================================================
# SLIDE 3: IMPACT QUANTIFICATION (BEFORE vs AFTER)
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3)
add_title_bar(slide3, 'SLIDE 3: Impact Quantification', 'Before vs After Metro Phase 2')

# Chart: Demand redistribution pies
add_chart_image(slide3, 's2_02_demand_redistribution.png', 0.3, 1.3, 6.5)

# Before/After table
add_table(slide3, 7, 1.3, 6, 2.8,
    ['Metric', 'Pre-Metro (H1)', 'Post-Metro (Q3)', 'Change'],
    [
        ['Daily Demand', '31,172', '31,566', '+1.3%'],
        ['Express Share', '24.2%', '31.0%', '+6.8pp'],
        ['Feeder Share', '24.4%', '18.5%', '-5.9pp'],
        ['City Share', '38.5%', '38.9%', '+0.4pp'],
        ['Congestion', '2.78', '3.09', '+10.9%'],
        ['Bus Speed', '31.1 km/h', '29.2 km/h', '-6.3%'],
        ['Cong. Elasticity', '+720 pax/lvl', '+1,447 pax/lvl', '+101%'],
        ['Pax/km (network)', '144.0', '146.1', '+1.5%'],
    ])

# Overload status
add_table(slide3, 7, 4.4, 6, 1.8,
    ['Route', 'Pre Pax/km', 'Post Pax/km', 'Change', 'Status'],
    [
        ['X28 (Exp)', '207', '267', '+28.9%', '🔴 CRITICAL'],
        ['X66 (Exp)', '123', '163', '+32.2%', '🔴 OVERLOADED'],
        ['X11 (Exp)', '113', '145', '+28.6%', '🔴 OVERLOADED'],
        ['F18 (Feed)', '150', '103', '-31.2%', '📉 Contracting'],
        ['F12 (Feed)', '88', '70', '-20.2%', '📉 At Risk'],
    ])

add_callout(slide3, 0.3, 6.4, 12.7, 0.7,
    '🔑 Total demand barely changed (+1.3%). But internal mix shifted dramatically: Express +7pp, Feeder -6pp. '
    'This is a REDISTRIBUTION, not contraction — same passengers, different routes.')

# ============================================================
# SLIDE 4: TRADE-OFF IDENTIFICATION
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4)
add_title_bar(slide4, 'SLIDE 4: Trade-Off Identification', 'Budget-Neutral Reallocation vs Fleet Expansion')

# Chart: Fleet revision
add_chart_image(slide4, 's2_05_fleet_revision.png', 0.3, 1.3, 6.5)

# Trade-off 1
add_rich_textbox(slide4, 7, 1.3, 6, 0.9, [
    ('Trade-Off 1: Express Capacity vs Feeder Viability', 12, True, RED),
    ('Every bus moved to Express improves 3,300 riders\' experience,', 10, False, DARK),
    ('but degrades 1,800 Feeder riders\' service.', 10, False, DARK),
])

# Trade-off 2
add_rich_textbox(slide4, 7, 2.3, 6, 0.9, [
    ('Trade-Off 2: Headway vs Coverage', 12, True, ORANGE),
    ('X28 headway: 20→11 min (38% better)', 10, False, GREEN),
    ('But C04: 20→35 min, F12: 20→44 min (barely viable)', 10, False, RED),
])

# Trade-off 3 as table
add_rich_textbox(slide4, 7, 3.3, 6, 0.5, [
    ('Trade-Off 3: Speed vs Sustainability', 12, True, DARK_BLUE),
])
add_table(slide4, 7, 3.9, 6, 1.5,
    ['Option', 'Speed', 'Cost', 'Duration'],
    [
        ['Reallocate only', 'Immediate', '₹0', '60-90 days'],
        ['Procure 8-10 buses', '60-90 days', '₹16-20 Cr', '12-18 months'],
        ['Hybrid (recommended)', 'Both', '₹16-20 Cr', 'Best of both'],
    ])

add_callout(slide4, 0.3, 5.7, 12.7, 0.8,
    '⚠️ RECOMMENDATION: Hybrid approach — reallocate 11 buses immediately as a bridge measure, '
    'while initiating procurement of 8-10 new buses to reach Express corridors by November 2025 (Winter Peak).')

# Revised fleet summary
add_table(slide4, 0.3, 6.6, 12.7, 0.8,
    ['X28', 'X11', 'X66', 'C03', 'C01', 'C02', 'C04', 'E22', 'E16', 'F25', 'F18', 'F12'],
    [
        ['10→11', '5→7', '6→8', '10→11', '7→7', '7→7', '4→4', '8→7', '4→4', '9→7', '7→5', '4→3'],
    ])

# ============================================================
# SLIDE 5: STAGE 3 OPTIMIZATION DIRECTION
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5)
add_title_bar(slide5, 'SLIDE 5: Stage 3 Optimization Direction', 'Not Final Answers — Strategic Direction')

# Chart
add_chart_image(slide5, 's2_06_shift_classification.png', 0.3, 1.3, 6.5)

# Three axes
add_rich_textbox(slide5, 7, 1.3, 6, 1.5, [
    ('Axis 1: Feeder-Metro Integration', 13, True, RED),
    ('• F18, F12 are contracting structurally — NOT recoverable', 10, False, DARK),
    ('• Direction: Restructure as metro-connector feeder routes', 10, False, DARK),
    ('• Question: Which stops should become metro feeder termini?', 10, False, GRAY),
])

add_rich_textbox(slide5, 7, 2.9, 6, 1.5, [
    ('Axis 2: Express Corridor Scaling', 13, True, ORANGE),
    ('• All 3 Express routes at OVER-CAPACITY post-reallocation', 10, False, DARK),
    ('• Direction: Bus-priority lanes + fleet expansion on X28', 10, False, DARK),
    ('• Question: Optimal Express fleet size under new demand regime?', 10, False, GRAY),
])

add_rich_textbox(slide5, 7, 4.5, 6, 1.5, [
    ('Axis 3: Dynamic Demand Management', 13, True, GREEN),
    ('• Congestion elasticity doubled — demand is more volatile', 10, False, DARK),
    ('• Direction: Dynamic headway triggers at congestion > Level 3', 10, False, DARK),
    ('• Question: How many buses as dynamic reserve?', 10, False, GRAY),
])

add_callout(slide5, 0.3, 6.1, 12.7, 1.0,
    '🎯 STAGE 3 THESIS: The challenge is not fighting the Metro shift — it\'s DESIGNING a bus network '
    'that complements the metro. Feeders must connect TO metro, not compete WITH it. '
    'Express must scale to absorb the redistributed demand. We are prepared to optimize this new equilibrium.')

# Footer on all slides
for slide in prs.slides:
    add_textbox(slide, 0.3, 7.15, 5, 0.3,
        'DECODE X 2026 | Stage 2 Structural Recalibration Brief', 8, False, GRAY)
    add_textbox(slide, 9, 7.15, 4, 0.3,
        f'Slide {prs.slides.index(slide) + 1} of 5', 8, False, GRAY, PP_ALIGN.RIGHT)

# ============================================================
# SAVE
# ============================================================
prs.save(str(OUT_FILE))
print(f"\n{'='*50}")
print(f"Presentation saved: {OUT_FILE}")
print(f"{'='*50}")
print(f"Size: {OUT_FILE.stat().st_size / 1024:.0f} KB")
print(f"\n⚠️  REMEMBER: Convert to PDF before submission!")
